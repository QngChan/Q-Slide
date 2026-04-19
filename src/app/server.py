from __future__ import annotations

import json
import time
import math
import os
import shutil
import subprocess
try:
    import pythoncom
except ImportError:
    pythoncom = None  # Linux/macOS — COM functions won't be used
import threading
import uuid
from pathlib import Path
from typing import Any

from fastapi import FastAPI, HTTPException, Request, WebSocket, WebSocketDisconnect
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles

ROOT_DIR = Path(__file__).resolve().parent.parent.parent
WEB_DIR = ROOT_DIR / "web"
MEDIA_DIR = ROOT_DIR / "data" / "media_uploads"
MEDIA_DIR.mkdir(parents=True, exist_ok=True)
PPT_SOURCE_DIR = MEDIA_DIR / "ppt_sources"
PPT_SOURCE_DIR.mkdir(parents=True, exist_ok=True)

ALLOWED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp"}
ALLOWED_VIDEO_EXTENSIONS = {".mp4", ".webm", ".mov", ".m4v", ".avi"}
ALLOWED_AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".m4a", ".aac", ".flac"}
ALLOWED_PPT_EXTENSIONS = {".pptx", ".docx", ".xlsx", ".xls", ".pdf"}


def media_kind_from_suffix(suffix: str) -> str | None:
    lowered = suffix.lower()
    if lowered in ALLOWED_IMAGE_EXTENSIONS:
        return "image"
    if lowered in ALLOWED_VIDEO_EXTENSIONS:
        return "video"
    if lowered in ALLOWED_AUDIO_EXTENSIONS:
        return "audio"
    if lowered == ".pdf":
        return "pdf"
    return None


def emu_to_px(value: int, slide_emu: int, target_px: int) -> int:
    if slide_emu <= 0:
        return 0
    return max(0, int((value / slide_emu) * target_px))


def paragraph_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text and p.text.strip()]
    return "\n".join(lines).strip()


def extract_pptx_slides(source_file: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation  # type: ignore[import]
    except Exception as exc:
        raise HTTPException(status_code=500, detail="Yerel sunum icin python-pptx gerekli.") from exc

    try:
        presentation = Presentation(str(source_file))
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PPTX açılamadi: {exc}") from exc

    slide_width_emu = int(presentation.slide_width)
    slide_height_emu = int(presentation.slide_height)
    if slide_width_emu <= 0 or slide_height_emu <= 0:
        raise HTTPException(status_code=500, detail="Geçersiz sunu boyutu.")

    target_width = 1920
    target_height = max(360, int(target_width * (slide_height_emu / slide_width_emu)))

    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        elements: list[dict[str, Any]] = []

        for shape in slide.shapes:
            left = emu_to_px(int(getattr(shape, "left", 0)), slide_width_emu, target_width)
            top = emu_to_px(int(getattr(shape, "top", 0)), slide_height_emu, target_height)
            width = emu_to_px(int(getattr(shape, "width", 0)), slide_width_emu, target_width)
            height = emu_to_px(int(getattr(shape, "height", 0)), slide_height_emu, target_height)
            if width <= 0 or height <= 0:
                continue

            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                image_id = uuid.uuid4().hex
                ext = f".{shape.image.ext.lower()}"
                if ext not in {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"}:
                    ext = ".png"
                image_name = f"{image_id}{ext}"
                image_path = MEDIA_DIR / image_name
                image_path.write_bytes(shape.image.blob)
                elements.append(
                    {
                        "type": "image",
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "src": f"/media/{image_name}",
                    }
                )
                continue

            text = paragraph_text(shape)
            if text:
                elements.append(
                    {
                        "type": "text",
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "text": text,
                        "font_size": 30,
                        "color": "#111827",
                        "align": "left",
                    }
                )

        slides.append(
            {
                "index": idx - 1,
                "width": target_width,
                "height": target_height,
                "background": "#ffffff",
                "elements": elements,
            }
        )

    if not slides:
        raise HTTPException(status_code=500, detail="Sunuda render edilecek slayt bulunamadi.")
    return slides


def _test_libreoffice_executable(path: str) -> bool:
    """Returns True if the executable actually responds to --version within 15 seconds."""
    try:
        result = subprocess.run(
            [path, "--headless", "--nologo", "--nodefault", "--norestore", "--version"],
            capture_output=True, text=True, timeout=15
        )
        if result.returncode != 0:
            return False
        output = (result.stdout + result.stderr).strip()
        return (not output) or ("libreoffice" in output.lower())
    except Exception:
        return False


def get_libreoffice_path() -> str | None:
    """Detects a *working* LibreOffice executable on Windows."""
    if os.name != "nt":
        p = shutil.which("soffice")
        if p and _test_libreoffice_executable(p):
            return p
        return None

    # Windows: try PATH first, then common install locations.
    for exe_name in ("soffice.exe", "soffice.com", "soffice"):
        p = shutil.which(exe_name)
        if p and _test_libreoffice_executable(p):
            return p

    bases = [
        os.environ.get("ProgramW6432"),
        os.environ.get("PROGRAMFILES"),
        os.environ.get("PROGRAMFILES(X86)"),
    ]
    for base in [b for b in bases if b]:
        candidates = [
            Path(base) / "LibreOffice" / "program" / "soffice.exe",
            Path(base) / "LibreOffice" / "program" / "soffice.com",
        ]
        for c in candidates:
            if c.exists() and _test_libreoffice_executable(str(c)):
                return str(c)

    return None


_installing_libreoffice = False


def download_and_install_libreoffice():
    global _installing_libreoffice
    _installing_libreoffice = True
    try:
        import urllib.request
        msi_url = "https://download.documentfoundation.org/libreoffice/stable/24.8.5/win/x86_64/LibreOffice_24.8.5_Win_x64.msi"
        msi_path = Path(os.environ.get("TEMP", ".")) / "LibreOffice_Setup.msi"

        # Download
        urllib.request.urlretrieve(msi_url, msi_path)

        # Install silently
        subprocess.run(["msiexec", "/i", str(msi_path.resolve()), "/qn", "/norestart"], check=True)
    except Exception as e:
        print(f"LibreOffice install error: {e}")
    finally:
        _installing_libreoffice = False


def _render_pdf_to_slides(pdf_path: Path, prefix: str = "pdf") -> list[dict[str, Any]]:
    """Shared helper: renders each page of a PDF as a 1920px-wide PNG slide."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("pymupdf not installed.")

    slides: list[dict[str, Any]] = []
    doc = fitz.open(str(pdf_path))
    try:
        target_width = 1920
        for i, page in enumerate(doc):
            rect = page.rect
            zoom = target_width / rect.width
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False, colorspace=fitz.csRGB)

            image_name = f"{prefix}_{uuid.uuid4().hex}.png"
            image_path = MEDIA_DIR / image_name
            pix.save(str(image_path.resolve()))

            slides.append({
                "index": i,
                "width": target_width,
                "height": int(rect.height * zoom),
                "background": "#ffffff",
                "elements": [{
                    "type": "image",
                    "left": 0,
                    "top": 0,
                    "width": target_width,
                    "height": int(rect.height * zoom),
                    "src": f"/media/{image_name}"
                }]
            })
    finally:
        doc.close()
    return slides


def extract_slides_via_libreoffice(source_file: Path) -> list[dict[str, Any]]:
    """Uses LibreOffice to convert PPTX/DOCX to PDF, then PyMuPDF to PNG."""
    lo_path = get_libreoffice_path()
    if not lo_path:
        raise RuntimeError("LibreOffice not found.")

    temp_dir = PPT_SOURCE_DIR / f"tmp_{uuid.uuid4().hex}"
    temp_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            [lo_path, "--headless", "--convert-to", "pdf", "--outdir", str(temp_dir.resolve()), str(source_file.resolve())],
            check=True,
            capture_output=True
        )

        pdf_file = next(temp_dir.glob("*.pdf"), None)
        if not pdf_file:
            raise RuntimeError("LibreOffice PDF conversion failed.")

        return _render_pdf_to_slides(pdf_file, "lo")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def extract_pptx_slides_as_images(source_file: Path) -> list[dict[str, Any]]:
    """Exports each slide as a PNG using Windows COM for perfect fidelity."""
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("pywin32 not installed")

    # Ensure COM is initialized for the current thread
    pythoncom.CoInitialize()

    ppt_app = None
    presentation = None
    try:
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
        # Ensure it's invisible to the user
        ppt_app.Visible = 1  # 1 is MsoTrue, but it often needs to be visible to export correctly
        # However, for automation we usually want it minimized or hidden if possible.
        # But in many environments, Background mode (Minimized) is safer.
        ppt_app.WindowState = 2  # 2 is ppWindowMinimized

        presentation = ppt_app.Presentations.Open(str(source_file.resolve()), ReadOnly=True, Untitled=False, WithWindow=False)

        slide_width = presentation.PageSetup.SlideWidth
        slide_height = presentation.PageSetup.SlideHeight
        # Standardize to 1920 width for consistency
        target_width = 1920
        target_height = int(target_width * (slide_height / slide_width))

        slides: list[dict[str, Any]] = []
        for i, slide in enumerate(presentation.Slides, start=1):
            image_name = f"ppt_{uuid.uuid4().hex}.png"
            image_path = MEDIA_DIR / image_name
            # Export slide as high res image
            slide.Export(str(image_path.resolve()), "PNG", ScaleX=target_width, ScaleY=target_height)

            slides.append({
                "index": i - 1,
                "width": target_width,
                "height": target_height,
                "background": "#000000",
                "elements": [{
                    "type": "image",
                    "left": 0,
                    "top": 0,
                    "width": target_width,
                    "height": target_height,
                    "src": f"/media/{image_name}"
                }]
            })
        return slides
    finally:
        if presentation:
            presentation.Close()
        if ppt_app:
            ppt_app.Quit()
        pythoncom.CoUninitialize()


def extract_docx_via_microsoft_word(source_file: Path) -> list[dict[str, Any]]:
    """Uses Microsoft Word COM to convert DOCX to PDF, then PyMuPDF to PNG."""
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("pywin32 not installed")

    pythoncom.CoInitialize()

    word_app = None
    doc = None
    temp_dir = PPT_SOURCE_DIR / f"tmp_{uuid.uuid4().hex}"
    temp_dir.mkdir(parents=True, exist_ok=True)
    try:
        word_app = win32com.client.DispatchEx("Word.Application")
        word_app.Visible = False
        word_app.DisplayAlerts = False

        doc = word_app.Documents.Open(str(source_file.resolve()), ReadOnly=True)
        pdf_path = temp_dir / f"{source_file.stem}.pdf"
        # wdFormatPDF = 17
        doc.SaveAs2(str(pdf_path.resolve()), FileFormat=17)
        doc.Close(False)
        doc = None

        return _render_pdf_to_slides(pdf_path, "word")
    finally:
        if doc:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word_app:
            try:
                word_app.Quit()
            except Exception:
                pass
        shutil.rmtree(temp_dir, ignore_errors=True)
        pythoncom.CoUninitialize()


def extract_xlsx_via_microsoft_excel(source_file: Path) -> list[dict[str, Any]]:
    """Uses Microsoft Excel COM to convert XLSX/XLS to PDF, then PyMuPDF to PNG."""
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("pywin32 not installed")

    pythoncom.CoInitialize()

    excel_app = None
    workbook = None
    temp_dir = PPT_SOURCE_DIR / f"tmp_{uuid.uuid4().hex}"
    temp_dir.mkdir(parents=True, exist_ok=True)
    try:
        excel_app = win32com.client.DispatchEx("Excel.Application")
        excel_app.Visible = False
        excel_app.DisplayAlerts = False

        workbook = excel_app.Workbooks.Open(str(source_file.resolve()), ReadOnly=True)
        pdf_path = temp_dir / f"{source_file.stem}.pdf"
        # xlTypePDF = 0
        workbook.ExportAsFixedFormat(0, str(pdf_path.resolve()))
        workbook.Close(False)
        workbook = None

        return _render_pdf_to_slides(pdf_path, "excel")
    finally:
        if workbook:
            try:
                workbook.Close(False)
            except Exception:
                pass
        if excel_app:
            try:
                excel_app.Quit()
            except Exception:
                pass
        shutil.rmtree(temp_dir, ignore_errors=True)
        pythoncom.CoUninitialize()


def _build_fallback_error_detail(suffix: str, errors: list[str]) -> str:
    """Builds a detailed Turkish user guidance message when all conversion methods fail."""
    format_names = {
        ".pptx": "PowerPoint (PPTX)",
        ".docx": "Word (DOCX)",
        ".xlsx": "Excel (XLSX)",
        ".xls": "Excel (XLS)",
    }
    fmt = format_names.get(suffix, suffix.upper())
    ms_app = {
        ".pptx": "Microsoft PowerPoint",
        ".docx": "Microsoft Word",
        ".xlsx": "Microsoft Excel",
        ".xls": "Microsoft Excel",
    }.get(suffix, "Microsoft Office")

    lines = [
        f"{fmt} dosyası açılamadı. Tüm dönüştürme yöntemleri başarısız oldu.",
        "",
        "Denenen yöntemler ve hatalar:",
    ]
    for i, err in enumerate(errors, 1):
        lines.append(f"  {i}. {err}")

    lines += [
        "",
        "Çözüm adımları:",
        f"  1) {ms_app} yüklüyse düzgün çalıştığından emin olun.",
        f"     Başlat menüsünden {ms_app} uygulamasını manuel açıp kapatmayı deneyin.",
        "  2) LibreOffice'i kurmak için Host panelindeki 'LibreOffice Kur' butonunu kullanın.",
        "     İndirme adresi: https://www.libreoffice.org/download",
        "  3) Dosyanın bozuk olmadığından emin olun —",
        f"     {ms_app} ile dosyayı manuel açıp kontrol edin.",
        "  4) Antivirüs/güvenlik yazılımı COM erişimini engelliyor olabilir.",
        "     Geçici olarak devre dışı bırakıp tekrar deneyin.",
        "  5) Sorun devam ederse uygulamayı yeniden başlatın.",
    ]
    return "\n".join(lines)


def extract_docx_slides(source_file: Path) -> list[dict[str, Any]]:
    try:
        from docx import Document  # type: ignore[import]
    except Exception as exc:
        raise HTTPException(status_code=500, detail="Yerel dokuman icin python-docx gerekli.") from exc

    try:
        doc = Document(str(source_file))
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DOCX açılamadı: {exc}") from exc

    target_width = 1920
    target_height = 1080  # Default height for docx "slides"

    slides: list[dict[str, Any]] = []
    current_elements: list[dict[str, Any]] = []
    current_top = 40

    def finalize_slide(idx: int):
        if not current_elements:
            return
        slides.append(
            {
                "index": idx,
                "width": target_width,
                "height": max(current_top + 40, target_height),
                "background": "#ffffff",
                "elements": list(current_elements),
            }
        )
        current_elements.clear()

    slide_idx = 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # If it's a heading, start a new slide (unless it's the first one)
        if para.style.name.startswith("Heading") and current_elements:
            finalize_slide(slide_idx)
            slide_idx += 1
            current_top = 40

        # Add text element
        current_elements.append(
            {
                "type": "text",
                "left": 100,
                "top": current_top,
                "width": 1720,
                "height": 60,
                "text": text,
                "font_size": 40 if para.style.name.startswith("Heading") else 30,
                "color": "#111827",
                "align": "left",
            }
        )
        current_top += 80 if para.style.name.startswith("Heading") else 50

        # Max elements per slide to prevent overflow
        if current_top > 900:
            finalize_slide(slide_idx)
            slide_idx += 1
            current_top = 40

    finalize_slide(slide_idx)

    if not slides:
        # Fallback for empty or unrecognized content
        slides.append({
            "index": 0,
            "width": target_width,
            "height": target_height,
            "background": "#ffffff",
            "elements": [{"type": "text", "left": 100, "top": 40, "width": 1720, "height": 60, "text": "İçerik bulunamadı.", "font_size": 30, "color": "#111827", "align": "center"}]
        })

    return slides


class ConferenceState:
    def __init__(self) -> None:
        self._clients: set[WebSocket] = set()
        self._media_items: list[dict[str, Any]] = []
        self._ppt_decks: list[dict[str, Any]] = []
        self._active_ppt_id: str | None = None
        self._current_ppt_slide: dict[str, Any] | None = None
        self._current_media_id: str | None = None
        self._bgm_media_id: str | None = None
        self._bgm_volume: float = 0.7
        now_ms = int(time.time() * 1000)
        self._video_control: dict[str, Any] = {"seq": 0, "action": "none", "at_ms": now_ms}
        self._bgm_control: dict[str, Any] = {"seq": 0, "action": "none", "volume": self._bgm_volume, "at_ms": now_ms}
        self._scroll_pct: float = 0.0
        self._video_pos: float = 0.0
        self._video_dur: float = 0.0
        self._bgm_pos: float = 0.0
        self._bgm_dur: float = 0.0

    async def get_state(self) -> dict[str, Any]:
        current = next((m for m in self._media_items if m["id"] == self._current_media_id), None)
        bgm = next((m for m in self._media_items if m["id"] == self._bgm_media_id), None)
        active_ppt = next((p for p in self._ppt_decks if p["id"] == self._active_ppt_id), None)
        return {
            "current_media": current,
            "current_ppt_slide": self._current_ppt_slide,
            "bgm_media": bgm,
            "active_presentation": {
                "id": active_ppt["id"],
                "name": active_ppt["name"],
                "slide_index": active_ppt["current_index"],
                "slide_count": len(active_ppt["slides"]),
            }
            if active_ppt
            else None,
            "bgm_volume": self._bgm_volume,
            "video_control": dict(self._video_control),
            "bgm_control": dict(self._bgm_control),
            "scroll_pct": self._scroll_pct,
            "video_pos": self._video_pos,
            "video_dur": self._video_dur,
            "bgm_pos": self._bgm_pos,
            "bgm_dur": self._bgm_dur,
            "updated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        }

    async def list_media(self) -> list[dict[str, Any]]:
        return list(self._media_items)

    async def add_media(self, item: dict[str, Any]) -> dict[str, Any]:
        self._media_items.append(item)
        return item

    async def update_media_bucket(self, media_id: str, bucket: str) -> dict[str, Any]:
        found = next((m for m in self._media_items if m["id"] == media_id), None)
        if not found:
            raise HTTPException(status_code=404, detail="Medya bulunamadı.")

        normalized_bucket = (bucket or "").strip().lower()
        if normalized_bucket not in {"presentation", "music"}:
            raise HTTPException(status_code=400, detail="Geçersiz medya alanı.")

        media_type = str(found.get("type") or "").strip().lower()
        if normalized_bucket == "presentation" and media_type not in {"image", "video", "pdf"}:
            raise HTTPException(status_code=400, detail="Bu medya sunum alanına taşınamaz.")
        if normalized_bucket == "music" and media_type not in {"audio", "video"}:
            raise HTTPException(status_code=400, detail="Bu medya müzik alanına taşınamaz.")

        found["bucket"] = normalized_bucket
        if normalized_bucket != "presentation" and self._current_media_id == media_id:
            self._current_media_id = None
        if normalized_bucket != "music" and self._bgm_media_id == media_id:
            self._bgm_media_id = None

        await self.broadcast_current()
        return found

    async def delete_media(self, media_id: str) -> dict[str, Any]:
        index = next((i for i, m in enumerate(self._media_items) if m["id"] == media_id), None)
        if index is None:
            raise HTTPException(status_code=404, detail="Medya bulunamadı.")

        removed = self._media_items.pop(index)
        if self._current_media_id == media_id:
            self._current_media_id = None
        if self._bgm_media_id == media_id:
            self._bgm_media_id = None

        media_name = Path(str(removed.get("url") or "")).name
        if media_name:
            file_path = MEDIA_DIR / media_name
            try:
                file_path.unlink(missing_ok=True)
            except Exception:
                pass

        await self.broadcast_current()
        return removed

    async def list_ppt_decks(self) -> dict[str, Any]:
        active = next((d for d in self._ppt_decks if d["id"] == self._active_ppt_id), None)
        return {
            "decks": [
                {
                    "id": d["id"],
                    "name": d["name"],
                    "slide_count": len(d["slides"]),
                    "current_index": d["current_index"],
                    "created_at": d["created_at"],
                    "preview_src": self._deck_preview_src(d),
                }
                for d in self._ppt_decks
            ],
            "active_deck_id": self._active_ppt_id,
            "active_slide_index": active["current_index"] if active else None,
        }

    @staticmethod
    def _deck_preview_src(deck: dict[str, Any]) -> str | None:
        slides = deck.get("slides") or []
        if not slides:
            return None
        first_slide = slides[0] or {}
        for element in first_slide.get("elements", []):
            if element.get("type") == "image" and element.get("src"):
                return str(element["src"])
        return None

    async def _activate_ppt_slide(self, deck: dict[str, Any]) -> None:
        if not deck["slides"]:
            self._current_ppt_slide = None
            await self.broadcast_current()
            return
        idx = max(0, min(len(deck["slides"]) - 1, int(deck["current_index"])))
        deck["current_index"] = idx
        self._current_media_id = None
        self._current_ppt_slide = dict(deck["slides"][idx])
        await self.broadcast_current()

    async def add_ppt_deck(self, deck: dict[str, Any]) -> dict[str, Any]:
        self._ppt_decks.append(deck)
        self._active_ppt_id = deck["id"]
        deck["current_index"] = 0
        await self._activate_ppt_slide(deck)
        return deck

    async def select_ppt_deck(self, deck_id: str) -> dict[str, Any]:
        deck = next((d for d in self._ppt_decks if d["id"] == deck_id), None)
        if not deck:
            raise HTTPException(status_code=404, detail="Sunu deck bulunamadı.")
        self._active_ppt_id = deck_id
        await self._activate_ppt_slide(deck)
        return deck

    async def step_active_ppt(self, delta: int) -> dict[str, Any]:
        if not self._active_ppt_id:
            raise HTTPException(status_code=400, detail="Aktif sunu bulunamadı.")
        deck = next((d for d in self._ppt_decks if d["id"] == self._active_ppt_id), None)
        if not deck:
            raise HTTPException(status_code=400, detail="Aktif sunu bulunamadı.")
        slides = deck["slides"]
        if not slides:
            raise HTTPException(status_code=400, detail="Aktif sunuda hiç sayfa yok.")
        deck["current_index"] = max(0, min(len(slides) - 1, int(deck["current_index"]) + delta))
        await self._activate_ppt_slide(deck)
        return {
            "deck_id": deck["id"],
            "slide_index": deck["current_index"],
            "slide_count": len(slides),
        }

    async def set_current_media(self, media_id: str) -> dict[str, Any]:
        found = next((m for m in self._media_items if m["id"] == media_id), None)
        if not found:
            raise HTTPException(status_code=404, detail="Medya bulunamadı.")
        if found["type"] == "audio":
            raise HTTPException(status_code=400, detail="Ses medyası sunucu tarafında görüntülenemez.")
        if found.get("bucket") != "presentation":
            raise HTTPException(status_code=400, detail="Sadece sunu medya sunucu tarafında görüntülenebilir.")
        self._current_media_id = media_id
        self._current_ppt_slide = None
        if found["type"] == "video":
            await self.emit_video_control("restart")
        await self.broadcast_current()
        return found

    async def set_bgm_media(self, media_id: str) -> dict[str, Any]:
        found = next((m for m in self._media_items if m["id"] == media_id), None)
        if not found:
            raise HTTPException(status_code=404, detail="Medya bulunamadı.")
        if found["type"] not in {"audio", "video"}:
            raise HTTPException(status_code=400, detail="Sadece ses/video medyası arka müzik olarak kullanılabilir.")
        if found.get("bucket") != "music":
            raise HTTPException(status_code=400, detail="Sadece müzik medyası arka müzik olarak kullanılabilir.")
        self._bgm_media_id = media_id
        await self.emit_bgm_control("restart")
        return found

    async def emit_video_control(self, action: str, time_seek: float | None = None) -> None:
        self._video_control = {
            "seq": int(self._video_control.get("seq", 0)) + 1,
            "action": action,
            "at_ms": int(time.time() * 1000),
        }
        if time_seek is not None:
            self._video_control["time"] = time_seek
        await self.broadcast_current()

    async def emit_bgm_control(self, action: str, volume: float | None = None, time_seek: float | None = None) -> None:
        if volume is not None:
            self._bgm_volume = max(0.0, min(1.0, float(volume)))
        self._bgm_control = {
            "seq": int(self._bgm_control.get("seq", 0)) + 1,
            "action": action,
            "volume": self._bgm_volume,
            "at_ms": int(time.time() * 1000),
        }
        if time_seek is not None:
            self._bgm_control["time"] = time_seek
        await self.broadcast_current()

    async def set_scroll_pct(self, pct: float) -> None:
        self._scroll_pct = max(0.0, min(1.0, float(pct)))
        await self.broadcast_current()

    async def sync_video_time(self, pos: float, dur: float) -> None:
        self._video_pos = float(pos)
        self._video_dur = float(dur)
        # We don't broadcast here to avoid flooding. Host will pick it up on polling or next state broadcast.

    async def sync_bgm_time(self, pos: float, dur: float) -> None:
        self._bgm_pos = float(pos)
        self._bgm_dur = float(dur)

    async def register(self, websocket: WebSocket) -> None:
        await websocket.accept()
        self._clients.add(websocket)
        await websocket.send_text(json.dumps({"type": "state", "payload": await self.get_state()}))

    async def unregister(self, websocket: WebSocket) -> None:
        self._clients.discard(websocket)

    async def broadcast(self, message: dict[str, Any]) -> None:
        encoded = json.dumps(message)
        stale: list[WebSocket] = []

        for client in list(self._clients):
            try:
                await client.send_text(encoded)
            except Exception:
                stale.append(client)

        for client in stale:
            self._clients.discard(client)

    async def broadcast_current(self) -> None:
        await self.broadcast({"type": "state", "payload": await self.get_state()})


conference_state = ConferenceState()
app = FastAPI(title="Q-Slide Konferans")
app.mount("/media", StaticFiles(directory=MEDIA_DIR), name="media")


@app.get("/")
async def viewer_page() -> FileResponse:
    return FileResponse(WEB_DIR / "viewer.html")


@app.get("/host")
async def host_page() -> FileResponse:
    return FileResponse(WEB_DIR / "host.html")


@app.get("/api/state")
async def get_state() -> JSONResponse:
    return JSONResponse(await conference_state.get_state())


@app.get("/api/media")
async def get_media() -> JSONResponse:
    return JSONResponse({"items": await conference_state.list_media()})


@app.get("/api/ppt")
async def get_ppt() -> JSONResponse:
    return JSONResponse(await conference_state.list_ppt_decks())


@app.post("/api/media/upload")
async def upload_media(request: Request, filename: str = "", bucket: str = "presentation") -> JSONResponse:
    original_name = (filename or request.headers.get("x-filename") or "file").strip()
    suffix = Path(original_name).suffix.lower()
    kind = media_kind_from_suffix(suffix)
    if not kind:
        raise HTTPException(status_code=400, detail="Sadece resim/video/ses ve PDF dosyaları desteklenir.")
    normalized_bucket = (bucket or "presentation").strip().lower()
    if normalized_bucket not in {"presentation", "music"}:
        raise HTTPException(status_code=400, detail="Müzik bucket sadece ses/video kabul eder.")
    if normalized_bucket == "presentation" and kind not in {"image", "video", "pdf"}:
        raise HTTPException(status_code=400, detail="Sunu bucket sadece resim/video/pdf kabul eder.")
    if normalized_bucket == "music" and kind not in {"audio", "video"}:
        raise HTTPException(status_code=400, detail="Müzik bucket sadece ses/video kabul eder.")

    media_id = uuid.uuid4().hex
    safe_name = f"{media_id}{suffix}"
    destination = MEDIA_DIR / safe_name
    raw = await request.body()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file content.")

    destination.write_bytes(raw)

    pages = []
    if kind == "pdf":
        try:
            slides = _render_pdf_to_slides(destination, "media_pdf")
            # Extract only the img src URLs from slides to store in media item
            pages = []
            for s in slides:
                if s.get("elements") and len(s["elements"]) > 0:
                    pages.append(s["elements"][0]["src"])
        except Exception as e:
            print(f"PDF media render error: {e}")
            # Fallback to empty pages, frontend will handle

    item = {
        "id": media_id,
        "name": original_name,
        "type": kind,
        "bucket": normalized_bucket,
        "url": f"/media/{safe_name}",
        "pages": pages,
        "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
    }
    saved = await conference_state.add_media(item)
    return JSONResponse(saved)


@app.post("/api/ppt/upload")
async def upload_ppt(request: Request, filename: str = "") -> JSONResponse:
    original_name = (filename or request.headers.get("x-filename") or "presentation.pptx").strip()
    suffix = Path(original_name).suffix.lower()
    if suffix not in ALLOWED_PPT_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Sadece .pptx/.docx/.xlsx/.xls dosyaları desteklenir.")

    raw = await request.body()
    if not raw:
        raise HTTPException(status_code=400, detail="Boş dosya içeriği.")

    deck_id = uuid.uuid4().hex
    source_file = PPT_SOURCE_DIR / f"{deck_id}{suffix}"
    source_file.write_bytes(raw)

    slides: list[dict[str, Any]] = []
    fallback_errors: list[str] = []

    if suffix in {".pptx", ".docx", ".xlsx", ".xls", ".pdf"}:
        # Strategy:
        # 1. Handle PDF directly if suffix is .pdf
        # 2. Try LibreOffice (Best for PPTX/DOCX/XLSX high-fidelity)
        # 3. Try Microsoft Office COM (PowerPoint / Word / Excel)
        # 4. python-pptx / python-docx text extraction (PPTX/DOCX only)
        # 5. Detailed user guidance error

        # --- 0. PDF Direct rendering ---
        if suffix == ".pdf":
            try:
                slides = _render_pdf_to_slides(source_file, "direct")
            except Exception as pdf_err:
                fallback_errors.append(f"PDF Render: {pdf_err}")
                print(f"PDF dönüşümü başarısız: {pdf_err}")

        # --- 1. LibreOffice ---
        lo_path = None
        if not slides:
            lo_path = get_libreoffice_path()
        if lo_path:
            try:
                slides = extract_slides_via_libreoffice(source_file)
            except Exception as lo_err:
                fallback_errors.append(f"LibreOffice: {lo_err}")
                print(f"LibreOffice dönüşümü başarısız: {lo_err}")
        else:
            fallback_errors.append("LibreOffice: Sistemde bulunamadı.")

        # --- 2. Microsoft Office COM ---
        if not slides:
            if suffix == ".pptx":
                try:
                    slides = extract_pptx_slides_as_images(source_file)
                except Exception as com_err:
                    fallback_errors.append(f"Microsoft PowerPoint COM: {com_err}")
                    print(f"PowerPoint COM dönüşümü başarısız: {com_err}")
            elif suffix == ".docx":
                try:
                    slides = extract_docx_via_microsoft_word(source_file)
                except Exception as word_err:
                    fallback_errors.append(f"Microsoft Word COM: {word_err}")
                    print(f"Word COM dönüşümü başarısız: {word_err}")
            elif suffix in {".xlsx", ".xls"}:
                try:
                    slides = extract_xlsx_via_microsoft_excel(source_file)
                except Exception as excel_err:
                    fallback_errors.append(f"Microsoft Excel COM: {excel_err}")
                    print(f"Excel COM dönüşümü başarısız: {excel_err}")

        # --- 3. python-pptx / python-docx text extraction ---
        if not slides:
            if suffix == ".pptx":
                try:
                    slides = extract_pptx_slides(source_file)
                except Exception as pptx_err:
                    fallback_errors.append(f"python-pptx: {pptx_err}")
            elif suffix == ".docx":
                try:
                    slides = extract_docx_slides(source_file)
                except Exception as docx_err:
                    fallback_errors.append(f"python-docx: {docx_err}")

        # --- 4. All methods failed → detailed user guidance ---
        if not slides:
            detail = _build_fallback_error_detail(suffix, fallback_errors)
            raise HTTPException(status_code=500, detail=detail)

    deck = {
        "id": deck_id,
        "name": original_name,
        "slides": slides,
        "current_index": 0,
        "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
    }
    created = await conference_state.add_ppt_deck(deck)
    return JSONResponse(
        {
            "id": created["id"],
            "name": created["name"],
            "slide_count": len(created["slides"]),
            "current_index": created["current_index"],
            "created_at": created["created_at"],
        }
    )


@app.post("/api/ppt/select")
async def select_ppt(payload: dict[str, Any]) -> JSONResponse:
    deck_id = str(payload.get("deck_id") or "").strip()
    if not deck_id:
        raise HTTPException(status_code=400, detail="deck_id gerekli.")
    selected = await conference_state.select_ppt_deck(deck_id)
    return JSONResponse(
        {
            "deck_id": selected["id"],
            "slide_index": selected["current_index"],
            "slide_count": len(selected["slides"]),
        }
    )


@app.post("/api/ppt/next")
async def next_ppt() -> JSONResponse:
    return JSONResponse(await conference_state.step_active_ppt(+1))


@app.post("/api/ppt/prev")
async def prev_ppt() -> JSONResponse:
    return JSONResponse(await conference_state.step_active_ppt(-1))


@app.post("/api/media/show")
async def show_media(payload: dict[str, Any]) -> JSONResponse:
    media_id = str(payload.get("media_id") or "").strip()
    if not media_id:
        raise HTTPException(status_code=400, detail="media_id gerekli.")
    selected = await conference_state.set_current_media(media_id)
    return JSONResponse(selected)


@app.post("/api/bgm/set")
async def set_bgm(payload: dict[str, Any]) -> JSONResponse:
    media_id = str(payload.get("media_id") or "").strip()
    if not media_id:
        raise HTTPException(status_code=400, detail="media_id gerekli.")
    selected = await conference_state.set_bgm_media(media_id)
    return JSONResponse(selected)


@app.post("/api/media/move")
async def move_media(payload: dict[str, Any]) -> JSONResponse:
    media_id = str(payload.get("media_id") or "").strip()
    bucket = str(payload.get("bucket") or "").strip().lower()
    if not media_id:
        raise HTTPException(status_code=400, detail="media_id gerekli.")
    if not bucket:
        raise HTTPException(status_code=400, detail="bucket gerekli.")
    moved = await conference_state.update_media_bucket(media_id, bucket)
    return JSONResponse(moved)


@app.post("/api/media/delete")
async def remove_media(payload: dict[str, Any]) -> JSONResponse:
    media_id = str(payload.get("media_id") or "").strip()
    if not media_id:
        raise HTTPException(status_code=400, detail="media_id gerekli.")
    removed = await conference_state.delete_media(media_id)
    return JSONResponse({"status": "ok", "removed": removed})


@app.post("/api/bgm/control")
async def control_bgm(payload: dict[str, Any]) -> JSONResponse:
    action = str(payload.get("action") or "").strip().lower()
    if action not in {"play", "pause", "stop", "restart", "volume", "seek"}:
        raise HTTPException(status_code=400, detail="Desteklenmeyen bgm eylemi.")
    volume = payload.get("volume", None)
    time_seek = payload.get("time", None)
    await conference_state.emit_bgm_control(action, volume=volume, time_seek=time_seek)
    return JSONResponse({"status": "ok"})


@app.post("/api/video/control")
async def control_video(payload: dict[str, Any]) -> JSONResponse:
    action = str(payload.get("action") or "").strip().lower()
    if action not in {"play", "pause", "restart", "seek"}:
        raise HTTPException(status_code=400, detail="Desteklenmeyen video eylemi.")
    time_seek = payload.get("time", None)
    await conference_state.emit_video_control(action, time_seek=time_seek)
    return JSONResponse({"status": "ok"})


@app.post("/api/video/sync")
async def sync_video(payload: dict[str, Any]) -> JSONResponse:
    pos = float(payload.get("pos", 0))
    dur = float(payload.get("dur", 0))
    await conference_state.sync_video_time(pos, dur)
    return JSONResponse({"status": "ok"})


@app.post("/api/bgm/sync")
async def sync_bgm(payload: dict[str, Any]) -> JSONResponse:
    pos = float(payload.get("pos", 0))
    dur = float(payload.get("dur", 0))
    await conference_state.sync_bgm_time(pos, dur)
    return JSONResponse({"status": "ok"})


@app.post("/api/scroll/sync")
async def sync_scroll(payload: dict[str, Any]) -> JSONResponse:
    raw_pct = payload.get("pct")
    if raw_pct is None:
        return JSONResponse({"status": "error", "message": "Missing pct"}, status_code=400)
    try:
        pct = float(raw_pct)
        if math.isnan(pct) or math.isinf(pct):
             return JSONResponse({"status": "error", "message": "Invalid pct value"}, status_code=400)
        print(f"Scroll sync: {pct:.2%}")
        await conference_state.set_scroll_pct(pct)
    except (ValueError, TypeError):
        return JSONResponse({"status": "error", "message": "Non-numeric pct"}, status_code=400)
    return JSONResponse({"status": "ok"})


_lo_cache: dict[str, Any] = {"path": None, "checked_at": 0.0}


@app.get("/api/system/status")
async def system_status() -> JSONResponse:
    now = time.time()
    # Re-check every 60 seconds or if currently installing
    if now - _lo_cache["checked_at"] > 60 or _installing_libreoffice:
        _lo_cache["path"] = get_libreoffice_path()
        _lo_cache["checked_at"] = now

    return JSONResponse({
        "libreoffice_found": _lo_cache["path"] is not None,
        "installing": _installing_libreoffice
    })


@app.post("/api/system/install_libreoffice")
async def install_libreoffice() -> JSONResponse:
    global _installing_libreoffice
    if _installing_libreoffice:
        return JSONResponse({"status": "zaten_installing"})
    
    if get_libreoffice_path():
        return JSONResponse({"status": "zaten_installed"})

    threading.Thread(target=download_and_install_libreoffice, daemon=True).start()
    return JSONResponse({"status": "started"})


try:
    from src.app.host import register_host_alive, unregister_host_alive
except ImportError:
    try:
        from .host import register_host_alive, unregister_host_alive
    except ImportError:
        def register_host_alive(client_id: str) -> None:
            pass
        def unregister_host_alive(client_id: str) -> None:
            pass


@app.post("/api/host/alive")
async def host_alive(payload: dict[str, Any]) -> JSONResponse:
    client_id = str(payload.get("client_id") or "").strip()
    if not client_id:
        raise HTTPException(status_code=400, detail="client_id gerekli.")
    register_host_alive(client_id)
    return JSONResponse({"status": "ok"})


@app.post("/api/host/bye")
async def host_bye(payload: dict[str, Any]) -> JSONResponse:
    client_id = str(payload.get("client_id") or "").strip()
    if client_id:
        unregister_host_alive(client_id)
    return JSONResponse({"status": "ok"})


@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket) -> None:
    await conference_state.register(websocket)
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        await conference_state.unregister(websocket)
